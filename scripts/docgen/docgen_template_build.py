"""
Build or modify a .docx or .pptx document template programmatically.

Supports:
  - Creating a new .docx from a JSON layout spec (tables, paragraphs, images)
  - Creating a new .pptx from a slide-based layout spec (`"format": "pptx"`)
  - Replacing tokens in an existing .docx (useful for branding updates)
  - Adding static images to an existing .docx
  - Listing/auditing structure of an existing .docx

Requires: pip install python-docx  (plus python-pptx for `"format": "pptx"`)

Usage:
  python scripts/docgen/docgen_template_build.py create layout.json --output template.docx
  python scripts/docgen/docgen_template_build.py create deck.json --output template.pptx
  python scripts/docgen/docgen_template_build.py replace template.docx --tokens '{"OldToken": "NewToken"}'
  python scripts/docgen/docgen_template_build.py audit template.docx
  python scripts/docgen/docgen_template_build.py --example > layout.json
  python scripts/docgen/docgen_template_build.py --example-pptx > deck.json
"""
import argparse
import json
import re
import sys

try:
    from docx import Document
    from docx.shared import Inches, Pt, Cm, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.enum.table import WD_TABLE_ALIGNMENT
except ImportError:
    print(
        "ERROR: python-docx required. Install with: pip install python-docx",
        file=sys.stderr,
    )
    sys.exit(1)


LAYOUT_EXAMPLE = {
    "_comment": "Layout spec for docgen_template_build.py create",
    "page": {"margin_inches": 0.75},
    "elements": [
        {
            "type": "image",
            "path": "/path/to/logo.png",
            "width_inches": 2.5,
            "alignment": "left",
        },
        {"type": "spacer"},
        {
            "type": "heading",
            "text": "INVOICE",
            "level": 1,
            "alignment": "right",
            "bold": True,
        },
        {"type": "spacer"},
        {
            "type": "table",
            "columns": 2,
            "style": "Table Grid",
            "rows": [
                ["Invoice Number:", "{{InvoiceNumber}}"],
                ["Invoice Date:", "{{InvoiceDate}}"],
                ["Due Date:", "{{DueDate}}"],
            ],
        },
        {"type": "spacer"},
        {
            "type": "paragraph",
            "text": "Bill To: {{AccountName}}",
            "bold": False,
            "size_pt": 11,
        },
        {"type": "spacer"},
        {
            "_comment": "Repeating section for line items",
            "type": "table",
            "columns": 4,
            "header": ["Description", "Quantity", "Unit Price", "Amount"],
            "rows": [
                [
                    "{{#InvoiceLines}}",
                    "",
                    "",
                    "",
                ],
                [
                    "{{Description}}",
                    "{{Quantity}}",
                    "{{UnitPrice}}",
                    "{{TotalAmount}}",
                ],
                ["{{/InvoiceLines}}", "", "", ""],
            ],
        },
        {"type": "spacer"},
        {
            "type": "paragraph",
            "text": "Total: {{TotalAmount}}",
            "bold": True,
            "size_pt": 12,
            "alignment": "right",
        },
    ],
}


PPTX_LAYOUT_EXAMPLE = {
    "_comment": "Slide layout spec for docgen_template_build.py create (PowerPoint)",
    "format": "pptx",
    "slide_size": {"width_inches": 13.333, "height_inches": 7.5},
    "slides": [
        {
            "elements": [
                {
                    "type": "textbox",
                    "left": 0.8, "top": 2.4, "width": 11.7, "height": 1.2,
                    "lines": [
                        {"text": "Proposal for {{AccountName}}", "size_pt": 40, "bold": True},
                        {"text": "Quote {{QuoteNumber}}", "size_pt": 18, "color": "#666666"},
                    ],
                }
            ]
        },
        {
            "elements": [
                {
                    "type": "textbox",
                    "left": 0.8, "top": 0.5, "width": 11.7, "height": 0.9,
                    "lines": [{"text": "Investment summary", "size_pt": 30, "bold": True}],
                },
                {
                    "type": "table",
                    "left": 0.8, "top": 1.6, "width": 11.7, "height": 1.2,
                    "columns": 4,
                    "col_widths": [5.4, 1.6, 2.3, 2.4],
                    "header": ["Product", "Qty", "Unit price", "Amount"],
                    "rows": [
                        [
                            "{{#Line}}{{ProductName}}",
                            "{{Quantity}}",
                            "{{NetUnitPrice}}",
                            "{{NetTotalPrice}}{{/Line}}",
                        ]
                    ],
                },
            ]
        },
    ],
}


PPTX_LAYOUT_EXAMPLE = {
    "_comment": "Slide layout spec for docgen_template_build.py create (PowerPoint)",
    "format": "pptx",
    "slide_size": {"width_inches": 13.333, "height_inches": 7.5},
    "slides": [
        {
            "elements": [
                {
                    "type": "textbox",
                    "left": 0.8, "top": 2.4, "width": 11.7, "height": 1.2,
                    "lines": [
                        {"text": "Proposal for {{AccountName}}", "size_pt": 40, "bold": True},
                        {"text": "Quote {{QuoteNumber}}", "size_pt": 18, "color": "#666666"},
                    ],
                }
            ]
        },
        {
            "elements": [
                {
                    "type": "textbox",
                    "left": 0.8, "top": 0.5, "width": 11.7, "height": 0.9,
                    "lines": [{"text": "Investment summary", "size_pt": 30, "bold": True}],
                },
                {
                    "type": "table",
                    "left": 0.8, "top": 1.6, "width": 11.7, "height": 1.2,
                    "columns": 4,
                    "col_widths": [5.4, 1.6, 2.3, 2.4],
                    "header": ["Product", "Qty", "Unit price", "Amount"],
                    "rows": [
                        [
                            "{{#Line}}{{ProductName}}",
                            "{{Quantity}}",
                            "{{NetUnitPrice}}",
                            "{{NetTotalPrice}}{{/Line}}",
                        ]
                    ],
                },
            ]
        },
    ],
}


def apply_alignment(paragraph, alignment_str):
    alignments = {
        "left": WD_ALIGN_PARAGRAPH.LEFT,
        "center": WD_ALIGN_PARAGRAPH.CENTER,
        "right": WD_ALIGN_PARAGRAPH.RIGHT,
    }
    if alignment_str and alignment_str.lower() in alignments:
        paragraph.alignment = alignments[alignment_str.lower()]


def _docx_rgb(color):
    rgb = color.lstrip("#")
    return RGBColor(int(rgb[0:2], 16), int(rgb[2:4], 16), int(rgb[4:6], 16))


def _set_row_exact_height(row, height_pt):
    """Force a table row to an exact height (twips)."""
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement

    trPr = row._tr.get_or_add_trPr()
    for existing in trPr.findall(qn("w:trHeight")):
        trPr.remove(existing)
    trHeight = OxmlElement("w:trHeight")
    trHeight.set(qn("w:val"), str(int(height_pt * 20)))  # points → twips
    trHeight.set(qn("w:hRule"), "exact")
    trPr.append(trHeight)


def _set_table_full_width(table):
    """Stretch a table to the full content width so banner/rule edges align."""
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement

    tbl = table._tbl
    tblPr = tbl.tblPr if tbl.tblPr is not None else OxmlElement("w:tblPr")
    if tbl.tblPr is None:
        tbl.insert(0, tblPr)
    for existing in tblPr.findall(qn("w:tblW")):
        tblPr.remove(existing)
    tblW = OxmlElement("w:tblW")
    tblW.set(qn("w:w"), "5000")  # 100% of content area
    tblW.set(qn("w:type"), "pct")
    tblPr.append(tblW)
    # Prefer fixed layout so edges stay aligned with accent rows.
    for existing in tblPr.findall(qn("w:tblLayout")):
        tblPr.remove(existing)
    layout = OxmlElement("w:tblLayout")
    layout.set(qn("w:type"), "fixed")
    tblPr.append(layout)


def _set_table_cell_margins(table, top=0, left=0, bottom=0, right=0):
    """Set table-level cell margins in twips (0 = flush cells, no internal gap)."""
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement

    tbl = table._tbl
    tblPr = tbl.tblPr if tbl.tblPr is not None else OxmlElement("w:tblPr")
    if tbl.tblPr is None:
        tbl.insert(0, tblPr)
    for existing in tblPr.findall(qn("w:tblCellMar")):
        tblPr.remove(existing)
    mar = OxmlElement("w:tblCellMar")
    for edge, val in (("top", top), ("left", left), ("bottom", bottom), ("right", right)):
        el = OxmlElement(f"w:{edge}")
        el.set(qn("w:w"), str(val))
        el.set(qn("w:type"), "dxa")
        mar.append(el)
    tblPr.append(mar)
    # Kill any inter-cell spacing that would open a white seam between rows.
    for existing in tblPr.findall(qn("w:tblCellSpacing")):
        tblPr.remove(existing)
    spacing = OxmlElement("w:tblCellSpacing")
    spacing.set(qn("w:w"), "0")
    spacing.set(qn("w:type"), "dxa")
    tblPr.append(spacing)


def _cant_split_row(row):
    """Prevent Word from splitting this table row across pages (avoids orphan banner lines)."""
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement

    tr = row._tr
    trPr = tr.get_or_add_trPr()
    for existing in trPr.findall(qn("w:cantSplit")):
        trPr.remove(existing)
    trPr.append(OxmlElement("w:cantSplit"))


def _shade_cell(cell, hex_color):
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement

    tc_pr = cell._tc.get_or_add_tcPr()
    for existing in tc_pr.findall(qn("w:shd")):
        tc_pr.remove(existing)
    shd = OxmlElement("w:shd")
    shd.set(qn("w:fill"), hex_color.lstrip("#"))
    shd.set(qn("w:val"), "clear")
    tc_pr.append(shd)


def _set_cell_borders_docx(cell, color=None, size_eighths=4, edges=None):
    """Set or clear cell borders.

    edges: optional dict like {"left": {"color": "#0B162A", "size": 48}, ...}
    When provided, overrides the uniform color/size_eighths per edge.
    """
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement

    tc_pr = cell._tc.get_or_add_tcPr()
    for existing in tc_pr.findall(qn("w:tcBorders")):
        tc_pr.remove(existing)
    borders = OxmlElement("w:tcBorders")
    for edge in ("top", "left", "bottom", "right"):
        el = OxmlElement(f"w:{edge}")
        edge_spec = (edges or {}).get(edge)
        if edge_spec is None and color is None:
            el.set(qn("w:val"), "nil")
        else:
            el.set(qn("w:val"), "single")
            if edge_spec:
                el.set(qn("w:sz"), str(edge_spec.get("size", size_eighths)))
                el.set(qn("w:color"), str(edge_spec.get("color", color or "000000")).lstrip("#"))
            else:
                el.set(qn("w:sz"), str(size_eighths))
                el.set(qn("w:color"), color.lstrip("#"))
            el.set(qn("w:space"), "0")
        borders.append(el)
    tc_pr.append(borders)


def _style_run(run, spec, defaults=None):
    defaults = defaults or {}
    if spec.get("bold", defaults.get("bold")):
        run.bold = True
    if spec.get("italic", defaults.get("italic")):
        run.italic = True
    size = spec.get("size_pt", defaults.get("size_pt"))
    if size:
        run.font.size = Pt(size)
    font = spec.get("font", defaults.get("font"))
    if font:
        run.font.name = font
    color = spec.get("color", defaults.get("color"))
    if color:
        run.font.color.rgb = _docx_rgb(color)


def _add_runs(paragraph, spec, defaults=None):
    """Add text or multi-run content to a paragraph."""
    defaults = defaults or {}
    runs = spec.get("runs")
    if runs:
        for rspec in runs:
            run = paragraph.add_run(rspec.get("text", ""))
            _style_run(run, rspec, defaults)
    else:
        run = paragraph.add_run(spec.get("text", ""))
        _style_run(run, spec, defaults)


def _set_paragraph_spacing(paragraph, spec):
    pf = paragraph.paragraph_format
    if spec.get("space_before_pt") is not None:
        pf.space_before = Pt(spec["space_before_pt"])
    if spec.get("space_after_pt") is not None:
        pf.space_after = Pt(spec["space_after_pt"])
    if spec.get("line_spacing"):
        pf.line_spacing = spec["line_spacing"]


def _apply_paragraph_rules(paragraph, color="#C83803", thickness_pt=2.25, space_pt=1, top=False, bottom=True):
    """Attach coloured top and/or bottom borders to a paragraph (title sandwich rules)."""
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement

    pPr = paragraph._p.get_or_add_pPr()
    for child in list(pPr):
        if child.tag == qn("w:pBdr"):
            pPr.remove(child)
    pBdr = OxmlElement("w:pBdr")
    sz = str(int(thickness_pt * 8))
    space = str(int(space_pt))
    hex_color = color.lstrip("#")
    for edge, enabled in (("top", top), ("bottom", bottom)):
        if not enabled:
            continue
        el = OxmlElement(f"w:{edge}")
        el.set(qn("w:val"), "single")
        el.set(qn("w:sz"), sz)
        el.set(qn("w:space"), space)
        el.set(qn("w:color"), hex_color)
        pBdr.append(el)
    pPr.append(pBdr)


def _apply_bottom_rule(paragraph, color="#C83803", thickness_pt=2.25, space_pt=1):
    """Backward-compatible alias — bottom rule only."""
    _apply_paragraph_rules(
        paragraph, color=color, thickness_pt=thickness_pt, space_pt=space_pt, top=False, bottom=True
    )


def _add_horizontal_rule(doc, color="#C83803", thickness_pt=2.25, space_before=6, space_after=10):
    """Orange (or coloured) rule via bottom border on a collapsed paragraph."""
    p = doc.add_paragraph()
    p.paragraph_format.space_before = Pt(space_before)
    p.paragraph_format.space_after = Pt(space_after)
    p.paragraph_format.line_spacing = Pt(1)
    run = p.add_run("")
    run.font.size = Pt(1)
    _apply_bottom_rule(p, color=color, thickness_pt=thickness_pt, space_pt=0)
    return p


def _set_page_background(doc, hex_color):
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement

    background = OxmlElement("w:background")
    background.set(qn("w:color"), hex_color.lstrip("#"))
    doc.element.insert(0, background)
    # Make Word/DocGen render the background shape.
    settings = doc.settings.element
    display = OxmlElement("w:displayBackgroundShape")
    display.set(qn("w:val"), "true")
    settings.append(display)


def _clear_cell(cell):
    for p in cell.paragraphs:
        p.clear()


def _write_page_number_cell(cell, color="#888888", font="Calibri", size_pt=8):
    """Write 'Page N of M' using Word PAGE/NUMPAGES fields (matches Bears preview footer)."""
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement

    _clear_cell(cell)
    para = cell.paragraphs[0]
    apply_alignment(para, "right")

    def _run(text):
        r = para.add_run(text)
        _style_run(r, {"size_pt": size_pt, "color": color, "font": font})
        return r

    def _field(instr):
        # Complex field: begin / instrText / separate / end
        def _fld_char(kind):
            r = OxmlElement("w:r")
            rPr = OxmlElement("w:rPr")
            sz = OxmlElement("w:sz")
            sz.set(qn("w:val"), str(int(size_pt * 2)))
            rPr.append(sz)
            color_el = OxmlElement("w:color")
            color_el.set(qn("w:val"), color.lstrip("#"))
            rPr.append(color_el)
            rFonts = OxmlElement("w:rFonts")
            rFonts.set(qn("w:ascii"), font)
            rFonts.set(qn("w:hAnsi"), font)
            rPr.append(rFonts)
            r.append(rPr)
            fc = OxmlElement("w:fldChar")
            fc.set(qn("w:fldCharType"), kind)
            r.append(fc)
            para._p.append(r)

        _fld_char("begin")
        r = OxmlElement("w:r")
        rPr = OxmlElement("w:rPr")
        sz = OxmlElement("w:sz")
        sz.set(qn("w:val"), str(int(size_pt * 2)))
        rPr.append(sz)
        color_el = OxmlElement("w:color")
        color_el.set(qn("w:val"), color.lstrip("#"))
        rPr.append(color_el)
        r.append(rPr)
        it = OxmlElement("w:instrText")
        it.set(qn("xml:space"), "preserve")
        it.text = f" {instr} "
        r.append(it)
        para._p.append(r)
        _fld_char("separate")
        # Placeholder text until Word/DocGen updates fields
        ph = para.add_run("1" if instr == "PAGE" else "13")
        _style_run(ph, {"size_pt": size_pt, "color": color, "font": font})
        _fld_char("end")

    _run("Page ")
    _field("PAGE")
    _run(" of ")
    _field("NUMPAGES")


def _write_cell(cell, text_or_lines, defaults=None, align=None):
    """Write a string, line-spec dict, or list of line-specs into a cell."""
    defaults = defaults or {}
    _clear_cell(cell)

    if isinstance(text_or_lines, dict):
        lines = [text_or_lines]
    elif isinstance(text_or_lines, list):
        lines = text_or_lines if text_or_lines else [{"text": ""}]
    else:
        lines = [{"text": "" if text_or_lines is None else str(text_or_lines)}]

    for i, line in enumerate(lines):
        if not isinstance(line, dict):
            line = {"text": str(line)}
        para = cell.paragraphs[0] if i == 0 else cell.add_paragraph()
        if align or line.get("align"):
            apply_alignment(para, line.get("align", align))
        _add_runs(para, line, defaults)
        _set_paragraph_spacing(para, line)


def create_from_layout(layout, output_path):
    doc = Document()
    page = layout.get("page", {})
    theme = layout.get("theme", {})
    defaults = {
        "font": page.get("font", "Calibri"),
        "size_pt": page.get("default_size_pt", 11),
        "color": page.get("default_color", theme.get("body", "#CCCCCC")),
    }

    margin = page.get("margin_inches", 0.75)
    for section in doc.sections:
        section.top_margin = Inches(page.get("margin_top_inches", margin))
        section.bottom_margin = Inches(page.get("margin_bottom_inches", margin))
        section.left_margin = Inches(page.get("margin_left_inches", margin))
        section.right_margin = Inches(page.get("margin_right_inches", margin))
        if page.get("different_first_page_header"):
            section.different_first_page_header_footer = True

    if page.get("background"):
        _set_page_background(doc, page["background"])

    # Content-page header (skipped on first page when different_first_page_header).
    header_spec = page.get("header")
    if header_spec:
        for section in doc.sections:
            header = section.header
            header.is_linked_to_previous = False
            # Clear default empty para
            for p in list(header.paragraphs):
                p.clear()
            table = header.add_table(rows=1, cols=2, width=Inches(7.5))
            table.autofit = True
            left, right = table.rows[0].cells
            _set_cell_borders_docx(left, None)
            _set_cell_borders_docx(right, None)
            # Left: logo + brand
            _clear_cell(left)
            lp = left.paragraphs[0]
            if header_spec.get("logo"):
                run = lp.add_run()
                try:
                    run.add_picture(header_spec["logo"], width=Inches(header_spec.get("logo_width_inches", 0.35)))
                except FileNotFoundError:
                    run.add_text("[logo]")
                lp.add_run("  ")
            brand = lp.add_run(header_spec.get("brand_text", "CHICAGO BEARS"))
            _style_run(brand, {
                "bold": True,
                "size_pt": header_spec.get("brand_size_pt", 9),
                "color": header_spec.get("brand_color", "#AAAAAA"),
                "font": defaults["font"],
            })
            # Right: subtitle
            _clear_cell(right)
            rp = right.paragraphs[0]
            apply_alignment(rp, "right")
            rr = rp.add_run(header_spec.get("right_text", ""))
            _style_run(rr, {
                "size_pt": header_spec.get("right_size_pt", 9),
                "color": header_spec.get("right_color", "#888888"),
                "font": defaults["font"],
            })
            # Optional rule under header. Set rule_color to null to omit — section
            # titles then own the underline without a second line above them.
            rule_color = header_spec.get("rule_color", "C83803")
            if rule_color:
                hp = header.add_paragraph()
                hp.paragraph_format.space_before = Pt(4)
                hp.paragraph_format.space_after = Pt(2)
                from docx.oxml.ns import qn
                from docx.oxml import OxmlElement
                pPr = hp._p.get_or_add_pPr()
                pBdr = OxmlElement("w:pBdr")
                bottom = OxmlElement("w:bottom")
                bottom.set(qn("w:val"), "single")
                bottom.set(qn("w:sz"), "18")
                bottom.set(qn("w:space"), "1")
                bottom.set(qn("w:color"), str(rule_color).lstrip("#"))
                pBdr.append(bottom)
                pPr.append(pBdr)

            if page.get("different_first_page_header"):
                # Leave first-page header empty.
                fh = section.first_page_header
                fh.is_linked_to_previous = False
                for p in fh.paragraphs:
                    p.clear()

    footer_spec = page.get("footer")
    if footer_spec:
        for section in doc.sections:
            footer = section.footer
            footer.is_linked_to_previous = False
            for p in footer.paragraphs:
                p.clear()
            # Thin rule
            rp = footer.add_paragraph()
            rp.paragraph_format.space_before = Pt(0)
            rp.paragraph_format.space_after = Pt(4)
            from docx.oxml.ns import qn
            from docx.oxml import OxmlElement
            pPr = rp._p.get_or_add_pPr()
            pBdr = OxmlElement("w:pBdr")
            top = OxmlElement("w:top")
            top.set(qn("w:val"), "single")
            top.set(qn("w:sz"), "6")
            top.set(qn("w:space"), "1")
            top.set(qn("w:color"), footer_spec.get("rule_color", "666666").lstrip("#"))
            pBdr.append(top)
            pPr.append(pBdr)

            table = footer.add_table(rows=1, cols=2, width=Inches(7.5))
            left, right = table.rows[0].cells
            _set_cell_borders_docx(left, None)
            _set_cell_borders_docx(right, None)
            _write_cell(left, footer_spec.get("left_text", ""), {
                "size_pt": 8, "color": footer_spec.get("color", "#888888"), "font": defaults["font"]
            })
            right_text = footer_spec.get("right_text", "")
            if right_text == "PAGE_NUMBER":
                _write_page_number_cell(
                    right,
                    color=footer_spec.get("color", "#888888"),
                    font=defaults["font"],
                    size_pt=8,
                )
            else:
                _write_cell(right, right_text, {
                    "size_pt": 8, "color": footer_spec.get("color", "#888888"), "font": defaults["font"]
                }, align="right")
            if footer_spec.get("disclaimer"):
                dp = footer.add_paragraph()
                apply_alignment(dp, "left")
                dr = dp.add_run(footer_spec["disclaimer"])
                _style_run(dr, {"size_pt": 7, "color": footer_spec.get("color", "#666666"), "font": defaults["font"]})

            if page.get("different_first_page_header"):
                ff = section.first_page_footer
                ff.is_linked_to_previous = False
                for p in ff.paragraphs:
                    p.clear()
                # Cover keeps a simple centered disclaimer in the body; first footer empty or minimal.
                if footer_spec.get("first_page_disclaimer"):
                    dp = ff.add_paragraph()
                    apply_alignment(dp, "center")
                    dr = dp.add_run(footer_spec["first_page_disclaimer"])
                    _style_run(dr, {"size_pt": 7, "color": "#666666", "font": defaults["font"]})

    for elem in layout.get("elements", []):
        elem_type = elem.get("type", "")

        if elem_type == "spacer":
            p = doc.add_paragraph("")
            if elem.get("space_after_pt") is not None:
                p.paragraph_format.space_after = Pt(elem["space_after_pt"])

        elif elem_type == "horizontal_rule":
            _add_horizontal_rule(
                doc,
                color=elem.get("color", theme.get("orange", "#C83803")),
                thickness_pt=elem.get("thickness_pt", 2.25),
                space_before=elem.get("space_before_pt", 4),
                space_after=elem.get("space_after_pt", 10),
            )

        elif elem_type in ("paragraph", "heading"):
            if elem_type == "heading":
                # Use a normal paragraph so we fully control colour on dark backgrounds
                # (Word's built-in Heading styles fight dark themes).
                p = doc.add_paragraph()
                size = elem.get("size_pt", {1: 22, 2: 16, 3: 13}.get(elem.get("level", 1), 16))
                hdefaults = {
                    **defaults,
                    "bold": True,
                    "size_pt": size,
                    "color": elem.get("color", theme.get("heading", "#8FA3B8")),
                }
            else:
                p = doc.add_paragraph()
                hdefaults = defaults
            apply_alignment(p, elem.get("alignment"))
            _add_runs(p, elem, hdefaults if elem_type == "heading" else {**defaults, **{k: elem[k] for k in ("bold", "italic", "size_pt", "color", "font") if k in elem}})
            _set_paragraph_spacing(p, elem)
            if elem.get("space_after_pt") is None and elem_type == "heading":
                p.paragraph_format.space_after = Pt(4)
            if elem.get("space_before_pt") is None and elem_type == "heading":
                p.paragraph_format.space_before = Pt(2)
            # Orange sandwich rules on the heading (above + below by default).
            rule = elem.get("underline_rule")
            if rule:
                _apply_paragraph_rules(
                    p,
                    color=rule.get("color", theme.get("orange", "#C83803")),
                    thickness_pt=rule.get("thickness_pt", 2.25),
                    space_pt=rule.get("space_pt", 2),
                    top=rule.get("above", True),
                    bottom=rule.get("below", True),
                )

        elif elem_type == "image":
            p = doc.add_paragraph()
            run = p.add_run()
            try:
                width = Inches(elem["width_inches"]) if elem.get("width_inches") else None
                run.add_picture(elem["path"], width=width)
            except FileNotFoundError:
                run.add_text(f"[IMAGE NOT FOUND: {elem['path']}]")
            apply_alignment(p, elem.get("alignment"))
            _set_paragraph_spacing(p, elem)

        elif elem_type == "banner":
            # Full-width shaded band — cover header, Total investment, metric strips.
            # Multi-line `content` becomes one shaded row per line so DocGen PDF
            # conversion keeps navy fill behind every line (multi-para cells lose it).
            fill = elem.get("fill", theme.get("navy", "#0B162A"))
            border = elem.get("border")  # None = no borders

            if elem.get("rows"):
                row_specs = elem["rows"]  # [[cell, ...], ...]
            elif elem.get("columns", 1) > 1:
                row_specs = [elem.get("cells", [elem.get("text", "")] * elem["columns"])]
            else:
                content = elem.get("content", elem.get("text", ""))
                if (
                    isinstance(content, list)
                    and content
                    and isinstance(content[0], dict)
                    and "text" in content[0]
                ):
                    row_specs = [[line] for line in content]
                else:
                    row_specs = [[content]]

            cols = max((len(r) if isinstance(r, list) else 1) for r in row_specs) if row_specs else 1
            table = doc.add_table(rows=len(row_specs), cols=cols)
            table.autofit = True
            _set_table_full_width(table)
            if elem.get("accent_bottom") or elem.get("flush_cells"):
                # Zero cell margins so navy ↔ orange meet with no white seam.
                _set_table_cell_margins(table, top=40, left=0, bottom=40, right=0)
            align = elem.get("align") or ["center"] * cols
            for r_i, row_content in enumerate(row_specs):
                if not isinstance(row_content, list):
                    row_content = [row_content]
                _cant_split_row(table.rows[r_i])
                for i, cell in enumerate(table.rows[r_i].cells):
                    _shade_cell(cell, fill)
                    _set_cell_borders_docx(cell, border, size_eighths=elem.get("border_size", 4))
                    content = row_content[i] if i < len(row_content) else ""
                    if isinstance(align, list):
                        align_i = align[i] if i < len(align) else align[0]
                    else:
                        align_i = align
                    cell_defaults = {
                        **defaults,
                        "color": elem.get("text_color", "#FFFFFF"),
                        "size_pt": elem.get("size_pt", 11),
                        "bold": elem.get("bold", False),
                    }
                    _write_cell(cell, content, cell_defaults, align=align_i)

            # Flush accent strip under the banner (same table = same left/right edges, no gap).
            accent = elem.get("accent_bottom")
            if accent:
                accent_color = accent.get("color", theme.get("orange", "#C83803"))
                thickness = accent.get("thickness_pt", 2.5)
                # Grow table by one row
                row = table.add_row()
                _cant_split_row(row)
                _set_row_exact_height(row, thickness)
                # Merge all cells so the accent is one continuous strip
                if cols > 1:
                    row.cells[0].merge(row.cells[cols - 1])
                cell = row.cells[0]
                _shade_cell(cell, accent_color)
                _set_cell_borders_docx(cell, None)
                _clear_cell(cell)
                # Collapse any leftover paragraph spacing in the accent cell
                for p in cell.paragraphs:
                    p.paragraph_format.space_before = Pt(0)
                    p.paragraph_format.space_after = Pt(0)
                    p.paragraph_format.line_spacing = Pt(1)

            if elem.get("space_after_pt", 8):
                sp = doc.add_paragraph("")
                sp.paragraph_format.space_after = Pt(elem.get("space_after_pt", 8))

        elif elem_type == "table":
            cols = elem.get("columns", 2)
            rows_data = elem.get("rows", [])
            header = elem.get("header")
            total_rows = len(rows_data) + (1 if header else 0)
            table = doc.add_table(rows=total_rows, cols=cols)
            table.style = elem.get("style", "Table Grid")
            border = elem.get("border", "333333")
            no_borders = elem.get("no_borders", False)

            row_offset = 0
            if header:
                for i, cell_text in enumerate(header[:cols]):
                    cell = table.rows[0].cells[i]
                    if elem.get("header_fill"):
                        _shade_cell(cell, elem["header_fill"])
                    if no_borders:
                        _set_cell_borders_docx(cell, None)
                    elif border:
                        _set_cell_borders_docx(cell, border, size_eighths=6)
                    _write_cell(
                        cell,
                        cell_text,
                        {
                            "bold": True,
                            "size_pt": elem.get("header_size_pt", 10),
                            "color": elem.get("header_text_color", "#FFFFFF"),
                            "font": defaults["font"],
                        },
                        align=(elem.get("align") or [None] * cols)[i] if elem.get("align") else None,
                    )
                row_offset = 1

            body_defaults = {
                "size_pt": elem.get("size_pt", defaults["size_pt"]),
                "color": elem.get("text_color", defaults["color"]),
                "font": defaults["font"],
            }
            for r, row_data in enumerate(rows_data):
                for c, cell_text in enumerate(row_data[:cols]):
                    cell = table.rows[r + row_offset].cells[c]
                    if elem.get("body_fill"):
                        _shade_cell(cell, elem["body_fill"])
                    if no_borders:
                        _set_cell_borders_docx(cell, None)
                    elif border:
                        _set_cell_borders_docx(
                            cell,
                            border,
                            size_eighths=elem.get("border_size", 4),
                            edges=elem.get("body_border_edges"),
                        )
                    _write_cell(
                        cell,
                        cell_text,
                        body_defaults,
                        align=(elem.get("align") or [None] * cols)[c] if elem.get("align") else None,
                    )

        elif elem_type == "page_break":
            doc.add_page_break()

    doc.save(output_path)
    print(f"Created: {output_path}")


def _pptx_rgb(color, theme=None):
    """Resolve a colour to RGBColor. `@name` looks up the layout's theme block."""
    from pptx.dml.color import RGBColor as PptxRGBColor

    if color.startswith("@"):
        if not theme or color[1:] not in theme:
            raise SystemExit(
                f"ERROR: colour '{color}' not found in the layout's theme block"
            )
        color = theme[color[1:]]
    rgb = color.lstrip("#")
    return PptxRGBColor(int(rgb[0:2], 16), int(rgb[2:4], 16), int(rgb[4:6], 16))


def _style_pptx_run(run, spec, theme=None):
    from pptx.util import Pt as PptxPt

    if spec.get("bold"):
        run.font.bold = True
    if spec.get("italic"):
        run.font.italic = True
    if spec.get("size_pt"):
        run.font.size = PptxPt(spec["size_pt"])
    if spec.get("font"):
        run.font.name = spec["font"]
    if spec.get("color"):
        run.font.color.rgb = _pptx_rgb(spec["color"], theme)
    # Letter spacing has no python-pptx accessor; `spc` is in 1/100 pt.
    if spec.get("letter_spacing_pt"):
        run.font._rPr.set("spc", str(int(spec["letter_spacing_pt"] * 100)))


def _pptx_align(name):
    from pptx.enum.text import PP_ALIGN

    return {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
    }.get((name or "left").lower(), PP_ALIGN.LEFT)


def _fill_text_frame(frame, spec, theme, default_size=14):
    """Write `lines` (or `text`) into a text frame with per-line styling."""
    from pptx.util import Pt as PptxPt
    from pptx.enum.text import MSO_ANCHOR

    frame.word_wrap = spec.get("word_wrap", True)
    for attr in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        if attr in spec:
            setattr(frame, attr, PptxPt(spec[attr]))
    if spec.get("anchor"):
        frame.vertical_anchor = {
            "top": MSO_ANCHOR.TOP,
            "middle": MSO_ANCHOR.MIDDLE,
            "bottom": MSO_ANCHOR.BOTTOM,
        }[spec["anchor"].lower()]

    lines = spec.get("lines") or [{"text": spec.get("text", "")}]
    for i, line in enumerate(lines):
        para = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        para.alignment = _pptx_align(line.get("align", spec.get("align")))
        if line.get("line_spacing"):
            para.line_spacing = line["line_spacing"]
        if line.get("space_before_pt"):
            para.space_before = PptxPt(line["space_before_pt"])
        if line.get("space_after_pt"):
            para.space_after = PptxPt(line["space_after_pt"])
        run = para.add_run()
        run.text = line.get("text", "")
        # Theme/layout default font applies unless the line overrides it.
        # Avoids the python-pptx default (Calibri) reading as generic AI type.
        merged = {"size_pt": default_size}
        if theme and theme.get("font"):
            merged["font"] = theme["font"]
        if spec.get("font"):
            merged["font"] = spec["font"]
        merged.update(line)
        _style_pptx_run(run, merged, theme)


def _set_cell_borders(cell, color=None, width_pt=1.0, theme=None):
    """Set or clear all four cell borders (python-pptx exposes no API for this)."""
    from lxml import etree

    ns = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    tc_pr = cell._tc.get_or_add_tcPr()
    # Order matters in the CT_TableCellProperties schema.
    for tag in ("lnL", "lnR", "lnT", "lnB"):
        for existing in tc_pr.findall(f"{ns}{tag}"):
            tc_pr.remove(existing)
        ln = etree.SubElement(tc_pr, f"{ns}{tag}")
        if color is None:
            etree.SubElement(ln, f"{ns}noFill")
        else:
            ln.set("w", str(int(width_pt * 12700)))
            fill = etree.SubElement(ln, f"{ns}solidFill")
            clr = etree.SubElement(fill, f"{ns}srgbClr")
            rgb = _pptx_rgb(color, theme)
            clr.set("val", f"{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}")


def create_pptx_from_layout(layout, output_path):
    """Build a tokenized .pptx from a slide-based layout spec.

    Every shape is explicitly positioned, because DocGen merges tokens into
    the text of existing shapes rather than reflowing a slide. A blank layout
    is used throughout so no placeholder prompt text survives into the
    rendered deck.
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches as PptxInches, Pt as PptxPt
    except ImportError:
        print(
            "ERROR: python-pptx required for \"format\": \"pptx\". "
            "Install with: pip install python-pptx",
            file=sys.stderr,
        )
        sys.exit(1)

    from pptx.enum.shapes import MSO_SHAPE

    SHAPES = {
        "rect": MSO_SHAPE.RECTANGLE,
        "rounded_rect": MSO_SHAPE.ROUNDED_RECTANGLE,
        "oval": MSO_SHAPE.OVAL,
    }

    prs = Presentation()
    size = layout.get("slide_size", {})
    prs.slide_width = PptxInches(size.get("width_inches", 13.333))
    prs.slide_height = PptxInches(size.get("height_inches", 7.5))
    theme = layout.get("theme", {})

    blank = prs.slide_layouts[6]

    for slide_spec in layout.get("slides", []):
        slide = prs.slides.add_slide(blank)

        # A full-bleed background rectangle must be added first so it sits behind
        # everything else; shape z-order follows insertion order.
        bg = slide_spec.get("background")
        if bg:
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = _pptx_rgb(bg, theme)
            shape.line.fill.background()
            shape.shadow.inherit = False

        for elem in slide_spec.get("elements", []):
            elem_type = elem.get("type", "")
            left = PptxInches(elem.get("left", 0.5))
            top = PptxInches(elem.get("top", 0.5))
            width = PptxInches(elem.get("width", 4.0))
            height = PptxInches(elem.get("height", 1.0))

            if elem_type == "textbox":
                frame = slide.shapes.add_textbox(left, top, width, height).text_frame
                _fill_text_frame(frame, elem, theme)

            elif elem_type == "shape":
                shape = slide.shapes.add_shape(
                    SHAPES[elem.get("shape", "rect")], left, top, width, height
                )
                if elem.get("fill"):
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = _pptx_rgb(elem["fill"], theme)
                else:
                    shape.fill.background()
                if elem.get("line"):
                    shape.line.color.rgb = _pptx_rgb(elem["line"], theme)
                    shape.line.width = PptxPt(elem.get("line_width_pt", 1))
                else:
                    shape.line.fill.background()
                shape.shadow.inherit = False
                if "corner_radius" in elem and shape.adjustments:
                    shape.adjustments[0] = elem["corner_radius"]
                if elem.get("lines") or elem.get("text"):
                    _fill_text_frame(shape.text_frame, elem, theme)

            elif elem_type == "table":
                cols = elem.get("columns", 2)
                rows_data = elem.get("rows", [])
                header = elem.get("header")
                total_rows = len(rows_data) + (1 if header else 0)

                table = slide.shapes.add_table(
                    total_rows, cols, left, top, width, height
                ).table

                # PowerPoint emphasises the first row and bands the body by
                # default; both fight explicit per-cell styling.
                table.first_row = elem.get("emphasise_first_row", False)
                table.horz_banding = elem.get("banded_rows", False)

                for i, w in enumerate((elem.get("col_widths") or [])[:cols]):
                    table.columns[i].width = PptxInches(w)
                for i, h in enumerate((elem.get("row_heights") or [])):
                    if i < total_rows:
                        table.rows[i].height = PptxInches(h)

                size_pt = elem.get("size_pt", 12)
                aligns = elem.get("align") or []
                border = elem.get("border")
                margins = elem.get("cell_margins_pt")

                def style_cell(cell, text, col, spec):
                    cell.text = ""
                    if spec.get("fill"):
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = _pptx_rgb(spec["fill"], theme)
                    else:
                        cell.fill.background()
                    if margins:
                        for attr, key in (
                            ("margin_left", "left"), ("margin_right", "right"),
                            ("margin_top", "top"), ("margin_bottom", "bottom"),
                        ):
                            if key in margins:
                                setattr(cell, attr, PptxPt(margins[key]))
                    _set_cell_borders(cell, border, elem.get("border_width_pt", 1.0), theme)
                    frame_spec = {
                        "lines": [{
                            "text": text,
                            "align": aligns[col] if col < len(aligns) else "left",
                            **{k: v for k, v in spec.items() if k != "fill"},
                        }],
                        "anchor": elem.get("cell_anchor", "middle"),
                    }
                    _fill_text_frame(cell.text_frame, frame_spec, theme, size_pt)

                row_offset = 0
                if header:
                    hspec = {
                        "bold": True,
                        "size_pt": elem.get("header_size_pt", size_pt),
                        **({"color": elem["header_text_color"]} if elem.get("header_text_color") else {}),
                        **({"fill": elem["header_fill"]} if elem.get("header_fill") else {}),
                    }
                    for i, text in enumerate(header[:cols]):
                        style_cell(table.cell(0, i), text, i, hspec)
                    row_offset = 1

                for r, row_data in enumerate(rows_data):
                    bspec = {
                        "size_pt": size_pt,
                        **({"color": elem["text_color"]} if elem.get("text_color") else {}),
                        **({"fill": elem["body_fill"]} if elem.get("body_fill") else {}),
                    }
                    for c, text in enumerate(row_data[:cols]):
                        style_cell(table.cell(r + row_offset, c), text, c, bspec)

            elif elem_type == "image":
                try:
                    slide.shapes.add_picture(elem["path"], left, top, width=width)
                except FileNotFoundError:
                    frame = slide.shapes.add_textbox(left, top, width, height).text_frame
                    frame.text = f"[IMAGE NOT FOUND: {elem['path']}]"

    prs.save(output_path)
    print(f"Created: {output_path}")


def replace_tokens(docx_path, token_map, output_path):
    doc = Document(docx_path)
    count = 0

    for para in doc.paragraphs:
        full_text = para.text
        for old, new in token_map.items():
            if old in full_text:
                for run in para.runs:
                    if old in run.text:
                        run.text = run.text.replace(old, new)
                        count += 1

    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for para in cell.paragraphs:
                    for old, new in token_map.items():
                        for run in para.runs:
                            if old in run.text:
                                run.text = run.text.replace(old, new)
                                count += 1

    doc.save(output_path or docx_path)
    print(f"Replaced {count} occurrence(s) in: {output_path or docx_path}")


def audit_docx(docx_path):
    doc = Document(docx_path)
    print(f"Auditing: {docx_path}")
    print(f"  Sections: {len(doc.sections)}")
    print(f"  Paragraphs: {len(doc.paragraphs)}")
    print(f"  Tables: {len(doc.tables)}")

    token_re = re.compile(r"\{\{([^}]+)\}\}")
    tokens_found = set()
    images_found = 0

    for para in doc.paragraphs:
        for match in token_re.finditer(para.text):
            tokens_found.add(match.group(1))
        for run in para.runs:
            if run.element.findall(
                ".//{http://schemas.openxmlformats.org/wordprocessingml/2006/main}drawing"
            ) or run.element.findall(
                ".//{http://schemas.openxmlformats.org/wordprocessingml/2006/main}pict"
            ):
                images_found += 1

    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for para in cell.paragraphs:
                    for match in token_re.finditer(para.text):
                        tokens_found.add(match.group(1))

    print(f"  Embedded images: {images_found}")
    print(f"  Unique tokens: {len(tokens_found)}")
    if tokens_found:
        fields = sorted(t for t in tokens_found if not t.startswith(("#", "/", "IMG_")))
        sections = sorted(t[1:] for t in tokens_found if t.startswith("#"))
        images = sorted(t for t in tokens_found if t.startswith("IMG_"))
        if fields:
            print(f"    Fields: {fields}")
        if sections:
            print(f"    Sections: {sections}")
        if images:
            print(f"    Images: {images}")


def main():
    parser = argparse.ArgumentParser(
        description="Build or modify .docx document templates"
    )
    subparsers = parser.add_subparsers(dest="command")

    create_p = subparsers.add_parser("create", help="Create .docx from layout spec")
    create_p.add_argument("layout", help="Path to layout JSON spec")
    create_p.add_argument("--output", "-o", required=True, help="Output .docx path")

    replace_p = subparsers.add_parser(
        "replace",
        help="Replace tokens in existing .docx (single-run tokens only)",
        epilog="Limitation: only replaces tokens contained within a single Word "
               "run. If Word splits a token across multiple runs (e.g. due to "
               "spell-check or formatting), the replacement will not match. "
               "Headers and footers are not searched. Use 'audit' to verify "
               "which tokens are visible.",
    )
    replace_p.add_argument("docx", help="Path to .docx file")
    replace_p.add_argument(
        "--tokens", required=True, help='JSON map: {"old": "new", ...}'
    )
    replace_p.add_argument("--output", "-o", help="Output path (default: overwrite)")

    audit_p = subparsers.add_parser("audit", help="Audit .docx structure")
    audit_p.add_argument("docx", help="Path to .docx file")

    parser.add_argument(
        "--example", action="store_true", help="Print example layout spec and exit"
    )
    parser.add_argument(
        "--example-pptx",
        action="store_true",
        help="Print example PowerPoint slide layout spec and exit",
    )

    args = parser.parse_args()

    if args.example:
        print(json.dumps(LAYOUT_EXAMPLE, indent=2))
        return

    if args.example_pptx:
        print(json.dumps(PPTX_LAYOUT_EXAMPLE, indent=2))
        return

    if not args.command:
        parser.print_help()
        return

    if args.command == "create":
        try:
            with open(args.layout) as f:
                layout = json.load(f)
        except (FileNotFoundError, json.JSONDecodeError) as e:
            print(f"ERROR: {e}", file=sys.stderr)
            sys.exit(1)

        is_pptx = layout.get("format") == "pptx" or "slides" in layout
        if is_pptx:
            create_pptx_from_layout(layout, args.output)
        else:
            create_from_layout(layout, args.output)

    elif args.command == "replace":
        try:
            token_map = json.loads(args.tokens)
        except json.JSONDecodeError as e:
            print(f"ERROR: Invalid JSON in --tokens: {e}", file=sys.stderr)
            sys.exit(1)
        replace_tokens(args.docx, token_map, args.output)

    elif args.command == "audit":
        audit_docx(args.docx)


if __name__ == "__main__":
    main()
